from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

# Title slide
title_slide = prs.slides.add_slide(prs.slide_layouts[0])
title_slide.shapes.title.text = "Python-PPTX Basics"
title_slide.placeholders[1].text = "Learn how to create and format slides with Python."

# Adding a custom text box
text_slide = prs.slides.add_slide(prs.slide_layouts[1])
text_box = text_slide.shapes.add_textbox(Inches(2), Inches(2), Inches(4), Inches(1.5))
text_box.text = "This is a custom text box created with python-pptx."

prs.save("output/04_add_text_boxes_example.pptx")
