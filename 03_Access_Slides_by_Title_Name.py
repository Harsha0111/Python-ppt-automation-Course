from pptx import Presentation

# Create a presentation
prs = Presentation()

# Function to get layout by name
def get_layout_by_name(prs, layout_name):
    for layout in prs.slide_layouts:
        if layout.name == layout_name:
            return layout
    return None

# Adding a slide using a named layout
title_slide_layout = get_layout_by_name(prs, "Title Slide")  # Assuming a layout with this name exists
content_slide_layout = get_layout_by_name(prs, "Title and Content")  # Assuming a layout with this name exists

# Create a title slide
if title_slide_layout:
    title_slide = prs.slides.add_slide(title_slide_layout)
    title_slide.shapes.title.text = "My Presentation"
    title_slide.placeholders[1].text = "Subtitle Here"

# Create a content slide
if content_slide_layout:
    content_slide = prs.slides.add_slide(content_slide_layout)
    content_slide.shapes.title.text = "Overview"
    content_slide.placeholders[1].text = "This slide contains the overview of the presentation."

# Save the presentation
prs.save("output/03_presentation_by_title_example.pptx")
