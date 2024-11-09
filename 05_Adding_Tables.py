from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Sales Data Table"

# Add table with 4 columns and 5 rows
rows, cols = 5, 4
table = slide.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(8), Inches(2)).table

# Column headers
headers = ["Product", "Units Sold", "Revenue", "Profit"]
for col, header in enumerate(headers):
    table.cell(0, col).text = header

# Sample data
data = [
    ["Product A", "500", "$10,000", "$4,000"],
    ["Product B", "300", "$6,000", "$2,400"],
    ["Product C", "700", "$14,000", "$5,600"],
]

# Fill table cells
for row, product_data in enumerate(data, start=1):
    for col, item in enumerate(product_data):
        table.cell(row, col).text = item

prs.save("output/05_table_example.pptx")
