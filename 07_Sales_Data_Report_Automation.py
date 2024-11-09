from pptx import Presentation
from pptx.util import Inches

# Initialize the presentation
prs = Presentation()

# Title Slide
title_slide = prs.slides.add_slide(prs.slide_layouts[0])
title_slide.shapes.title.text = "Monthly Sales Report"
title_slide.placeholders[1].text = "Month: August 2024"

# Summary Slide
summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
summary_slide.shapes.title.text = "Sales Summary"
summary_slide.placeholders[1].text = "This report provides an overview of the sales performance for August 2024."

# Sales Data Table Slide
table_slide = prs.slides.add_slide(prs.slide_layouts[5])
table_slide.shapes.title.text = "Sales Data"

# Define table structure and populate with data
rows, cols = 5, 4
table = table_slide.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(8), Inches(2)).table
headers = ["Product", "Units Sold", "Revenue", "Profit"]
data = [
    ["Product A", "500", "$10,000", "$4,000"],
    ["Product B", "300", "$6,000", "$2,400"],
    ["Product C", "700", "$14,000", "$5,600"],
]

# Fill headers
for col, header in enumerate(headers):
    table.cell(0, col).text = header

# Fill data rows
for row, product_data in enumerate(data, start=1):
    for col, item in enumerate(product_data):
        table.cell(row, col).text = item

# Conclusion Slide
conclusion_slide = prs.slides.add_slide(prs.slide_layouts[1])
conclusion_slide.shapes.title.text = "Conclusion"
conclusion_slide.placeholders[1].text = "Thank you for reviewing the sales report."

# Save the presentation
prs.save("output/07_Monthly_Sales_Report_August_2024.pptx")
