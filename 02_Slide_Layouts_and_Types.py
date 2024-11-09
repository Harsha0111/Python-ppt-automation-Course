from pptx import Presentation

prs = Presentation()

# Add title slide layout
title_slide = prs.slides.add_slide(prs.slide_layouts[0])
title_slide.shapes.title.text = "Welcome to Python-PPTX"
title_slide.placeholders[1].text = "Automate PowerPoint presentations with Python"

# Add content slide layout
content_slide = prs.slides.add_slide(prs.slide_layouts[1])
content_slide.shapes.title.text = "Agenda"
content_slide.placeholders[1].text = "1. Introduction\n2. Examples\n3. Project"

prs.save("output/02_slide_layouts_example.pptx")
