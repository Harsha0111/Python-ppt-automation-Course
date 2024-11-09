from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[1])

# Title and subtitle with formatting
title = slide.shapes.title
title.text = "Formatted Text Example"
title.text_frame.paragraphs[0].font.size = Pt(32)
title.text_frame.paragraphs[0].font.bold = True

# Custom text formatting
subtitle = slide.placeholders[1]
subtitle.text = "Customized font color and size."

# Changing font properties
paragraph = subtitle.text_frame.paragraphs[0]
paragraph.font.size = Pt(18)
paragraph.font.color.rgb = RGBColor(0, 102, 204)

prs.save("output/05_formatted_text_example.pptx")
